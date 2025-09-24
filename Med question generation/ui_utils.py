import io
from typing import List, Dict, Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

def create_powerpoint_from_slides(slides_data: List[Dict[str, Any]], presentation_title: str) -> bytes:
    """Generates a PowerPoint presentation from slide data with improved styling."""
    prs = Presentation()
    # Set slide dimensions to widescreen (standard 16:9 ratio; adjust if needed)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Title Slide ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    # Split "Book: Chapter" into a main title and a subtitle for better aesthetics
    if ":" in presentation_title:
        book, chapter = presentation_title.split(":", 1)
        title.text = chapter.strip()
        subtitle.text = f"Source: {book.strip()}"
    else:
        title.text = presentation_title

    # Style title
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # Style subtitle
    if subtitle.has_text_frame and subtitle.text:
        p = subtitle.text_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.italic = True
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # --- Content Slides ---
    content_slide_layout = prs.slide_layouts[1]  # Title and Content layout

    for slide_data in slides_data:
        slide = prs.slides.add_slide(content_slide_layout)
        
        # Add speaker notes if they exist
        if slide_data.get("speaker_notes"):
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data["speaker_notes"]

        # Set Title (left-aligned for a more professional look)
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get("title", "Untitled Slide")
        
        title_p = title_shape.text_frame.paragraphs[0]
        title_p.font.bold = True
        title_p.font.size = Pt(36)
        title_p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

        # Adjust title position for standard margins
        title_shape.left = Inches(0.5)
        title_shape.top = Inches(0.25)
        title_shape.width = Inches(12.333)
        title_shape.height = Inches(1.25)

        # Set Body Content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()  # Clear default placeholder text
        tf.word_wrap = True
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

        # Adjust body placeholder for better alignment and padding
        body_shape.left = Inches(0.75)
        body_shape.width = Inches(11.833)
        body_shape.top = Inches(1.75)
        body_shape.height = Inches(5.0)

        for bullet in slide_data.get("bullets", []):
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(24)  # Increased for readability
            p.space_before = Pt(12) # Added space between bullets
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    # --- Thank You Slide ---
    thank_you_layout = prs.slide_layouts[5]  # Title Only layout
    slide = prs.slides.add_slide(thank_you_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Thank You"
    
    # Center the text vertically and horizontally
    title_shape.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    title_shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # --- Add Footer and Slide Numbers to Content Slides ---
    num_content_slides = len(slides_data)
    for i, slide in enumerate(prs.slides):
        # Skip numbering and footer for the title and thank you slides
        if i == 0 or i > num_content_slides:
            continue

        # Add a footer with the chapter title for context
        footer_box = slide.shapes.add_textbox(
            left=Inches(0.5), top=Inches(6.75), width=Inches(10.0), height=Inches(0.2)
        )
        p = footer_box.text_frame.add_paragraph()
        if ":" in presentation_title:
            _, chapter = presentation_title.split(":", 1)
            p.text = chapter.strip()
        else:
            p.text = presentation_title
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(128, 128, 128)  # Subtle gray
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

        # Add slide number (e.g., "1 / 10")
        num_box = slide.shapes.add_textbox(
            left=Inches(12.0), top=Inches(6.75), width=Inches(1.0), height=Inches(0.2)
        )
        p = num_box.text_frame.add_paragraph()
        p.text = f"{i} / {num_content_slides}"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(128, 128, 128)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

    # Save the presentation to an in-memory stream
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream.getvalue()